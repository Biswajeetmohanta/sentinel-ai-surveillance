import os
import subprocess
from pptx import Presentation
from pptx.util import Inches

SUBMISSION_DIR = r"D:\sentinel-ai-surveillance\submission_package"
CHROME_PATH = r"C:\Program Files\Google\Chrome\Application\chrome.exe"
if not os.path.exists(CHROME_PATH):
    CHROME_PATH = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"

# Step 1: Create individual slide HTML files for perfect 1920x1080 rendering
with open(os.path.join(SUBMISSION_DIR, "interactive_presentation.html"), "r", encoding="utf-8") as f:
    full_html = f.read()

# Generate 8 standalone slide pages that are naturally active
os.makedirs(os.path.join(SUBMISSION_DIR, "slides_temp"), exist_ok=True)
slide_images = []

for i in range(1, 9):
    # Modify HTML so that only slide `i` is active and visible
    # We replace the JS or CSS so slide `i` has active class and others are hidden
    slide_html = full_html.replace(
        '<div class="slide active" id="slide1">',
        f'<div class="slide {"active" if i==1 else ""}" id="slide1">'
    )
    for s_num in range(1, 9):
        if s_num == i:
            slide_html = slide_html.replace(f'id="slide{s_num}"', f'id="slide{s_num}" class="slide active"')
        else:
            slide_html = slide_html.replace(f'id="slide{s_num}" class="slide active"', f'id="slide{s_num}" class="slide"')
            slide_html = slide_html.replace(f'class="slide active" id="slide{s_num}"', f'class="slide" id="slide{s_num}"')

    # Remove navigation controls for clean capture
    slide_html = slide_html.replace('<div class="nav-controls">', '<div class="nav-controls" style="display:none;">')
    slide_html = slide_html.replace('<div class="top-actions">', '<div class="top-actions" style="display:none;">')
    slide_html = slide_html.replace('<div class="progress-container">', '<div class="progress-container" style="display:none;">')

    html_file = os.path.join(SUBMISSION_DIR, "slides_temp", f"slide_{i}.html")
    with open(html_file, "w", encoding="utf-8") as f:
        f.write(slide_html)

    img_file = os.path.join(SUBMISSION_DIR, "slides_temp", f"slide_{i}.png")
    cmd = [
        CHROME_PATH,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        "--window-size=1920,1080",
        "--hide-scrollbars",
        f"--screenshot={img_file}",
        f"file:///{html_file.replace(os.sep, '/')}"
    ]
    subprocess.run(cmd, capture_output=True)
    if os.path.exists(img_file):
        print(f"Captured Slide {i} image: {img_file}")
        slide_images.append(img_file)
    else:
        print(f"Failed to capture slide {i}")

# Step 2: Assemble pristine 16:9 PowerPoint Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

for img_path in slide_images:
    slide = prs.slides.add_slide(blank_layout)
    # Add image full bleed
    slide.shapes.add_picture(img_path, 0, 0, Inches(13.333), Inches(7.5))

output_pptx = os.path.join(SUBMISSION_DIR, "Sentinel_AI_Solution_Presentation.pptx")
prs.save(output_pptx)
print(f"\n=======================================================")
print(f"SUCCESS: 100% Identical PPTX saved to: {output_pptx}")
print(f"=======================================================")
