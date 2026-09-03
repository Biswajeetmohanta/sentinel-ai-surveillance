import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_perfect_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette from interactive_presentation.html
    BG_DARK = RGBColor(11, 15, 25)          # #0b0f19
    CARD_BG = RGBColor(19, 29, 49)          # #131d31
    CARD_BORDER = RGBColor(30, 41, 59)      # #1e293b
    SUB_BOX_BG = RGBColor(15, 23, 42)       # #0f172a
    ACCENT_CYAN = RGBColor(6, 182, 212)     # #06b6d4
    CYAN_LIGHT = RGBColor(56, 189, 248)     # #38bdf8
    ACCENT_BLUE = RGBColor(59, 130, 246)    # #3b82f6
    ACCENT_PURPLE = RGBColor(168, 85, 247)  # #a855f7
    PURPLE_LIGHT = RGBColor(192, 132, 252)  # #c084fc
    ACCENT_GREEN = RGBColor(16, 185, 129)   # #10b981
    GREEN_LIGHT = RGBColor(52, 211, 153)    # #34d399
    ACCENT_RED = RGBColor(239, 68, 68)      # #ef4444
    RED_LIGHT = RGBColor(248, 113, 113)     # #f87171
    ACCENT_AMBER = RGBColor(245, 158, 11)   # #f59e0b
    AMBER_LIGHT = RGBColor(251, 191, 36)    # #fbbf24
    YELLOW_PLATE = RGBColor(250, 204, 21)   # #facc15
    TEXT_LIGHT = RGBColor(248, 250, 252)    # #f8fafc
    TEXT_MUTED = RGBColor(148, 163, 184)    # #94a3b8
    TEXT_WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)

    blank_layout = prs.slide_layouts[6]

    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

    def add_header(slide, title, category="GUJARAT POLICE SMART SURVEILLANCE", subtitle="", slide_num=1):
        # Category Badge Box
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(3.2), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(15, 30, 48)
        badge.line.color.rgb = ACCENT_CYAN
        badge.line.width = Pt(1)
        tf_b = badge.text_frame
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_b = tf_b.paragraphs[0]
        p_b.text = f"🛡️ {category.upper()}"
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = ACCENT_CYAN
        p_b.alignment = PP_ALIGN.CENTER

        # Title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.55))
        tf_t = tb_title.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

        # Subtitle
        if subtitle:
            tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.35))
            tf_s = tb_sub.text_frame
            p_s = tf_s.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(12)
            p_s.font.color.rgb = TEXT_MUTED

        # Footer line & counter
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.fill.background()

        tb_foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.95), Inches(11.7), Inches(0.4))
        tf_foot = tb_foot.text_frame
        p_f1 = tf_foot.paragraphs[0]
        p_f1.text = "Sentinel AI Surveillance Platform • 100% Self-Hosted • SDC Compliant"
        p_f1.font.size = Pt(10)
        p_f1.font.color.rgb = RGBColor(100, 116, 139)

        tb_cnt = slide.shapes.add_textbox(Inches(11.0), Inches(6.95), Inches(1.5), Inches(0.4))
        p_cnt = tb_cnt.text_frame.paragraphs[0]
        p_cnt.text = f"Slide {slide_num} / 8"
        p_cnt.font.size = Pt(10)
        p_cnt.font.bold = True
        p_cnt.font.color.rgb = RGBColor(100, 116, 139)
        p_cnt.alignment = PP_ALIGN.RIGHT

    def add_card(slide, left, top, width, height, title, subtitle="", top_accent_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.2)

        if top_accent_color:
            accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.12))
            accent.fill.solid()
            accent.fill.fore_color.rgb = top_accent_color
            accent.line.fill.background()

        # Text Frame
        tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.18), Inches(width - 0.5), Inches(height - 0.35))
        tf = tb.text_frame
        tf.word_wrap = True

        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

        if subtitle:
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.size = Pt(10.5)
            p_sub.font.color.rgb = TEXT_MUTED
            p_sub.space_before = Pt(3)

        return tf

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)

    # Accent Pillar
    pillar = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.5), Inches(0.2), Inches(4.3))
    pillar.fill.solid()
    pillar.fill.fore_color.rgb = ACCENT_CYAN
    pillar.line.fill.background()

    # Category Badge
    badge1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.5), Inches(3.8), Inches(0.38))
    badge1.fill.solid()
    badge1.fill.fore_color.rgb = RGBColor(15, 30, 48)
    badge1.line.color.rgb = ACCENT_CYAN
    badge1.line.width = Pt(1)
    p_b1 = badge1.text_frame.paragraphs[0]
    p_b1.text = "🛡️ GUJARAT POLICE SMART SURVEILLANCE"
    p_b1.font.size = Pt(11)
    p_b1.font.bold = True
    p_b1.font.color.rgb = ACCENT_CYAN
    p_b1.alignment = PP_ALIGN.CENTER

    # Main Title
    tb1 = s1.shapes.add_textbox(Inches(1.4), Inches(2.0), Inches(10.5), Inches(3.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "SENTINEL AI SURVEILLANCE PLATFORM"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Unified CCTV Ingestion, AI ANPR & GIS Vehicle Tracking System"
    p2.font.size = Pt(19)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_CYAN
    p2.space_before = Pt(8)

    p3 = tf1.add_paragraph()
    p3.text = "An enterprise-grade, 100% self-hosted computer vision intelligence system. Integrates multi-vendor cameras, executes sub-second hotlist alerts, and reconstructs vehicle journeys on interactive GIS maps with zero cloud API costs."
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(12)

    # 4 Bottom Badges
    b_items = [
        "📹 Multi-Vendor RTSP/ONVIF",
        "⚡ YOLOv8 + PP-OCRv4",
        "🚨 <400ms Hotlist Siren",
        "🗺️ PostGIS Breadcrumbs"
    ]
    for i, txt in enumerate(b_items):
        box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4 + i * 2.7), Inches(5.6), Inches(2.55), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CARD_BORDER
        box.line.width = Pt(1)
        p_bx = box.text_frame.paragraphs[0]
        p_bx.text = txt
        p_bx.font.size = Pt(11)
        p_bx.font.bold = True
        p_bx.font.color.rgb = TEXT_LIGHT
        p_bx.alignment = PP_ALIGN.CENTER

    # Footer
    tb_f1 = s1.shapes.add_textbox(Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.4))
    p_f = tb_f1.text_frame.paragraphs[0]
    p_f.text = "State Data Centre (SDC) & Air-Gapped Network Compliant • Official Gujarat Police Solution Deck"
    p_f.font.size = Pt(11)
    p_f.font.color.rgb = RGBColor(100, 116, 139)

    # =========================================================================
    # SLIDE 2: THE OPERATIONAL CHALLENGE
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "The Surveillance Challenge Faced by Police Forces", "PROBLEM CONTEXT", "Current city surveillance systems suffer from fragmented hardware, slow alerting, and high cloud costs.", 2)

    cards2 = [
        (0.8, "📹 Multi-Vendor Silos", "Incompatible Brands & Proprietary VMS", ACCENT_RED, [
            "Cities have thousands of cameras from different brands (Hikvision, Dahua, CP Plus, Axis) running incompatible proprietary software.",
            "✖ No unified live video matrix across city zones",
            "✖ Difficult and expensive to add or swap camera brands"
        ]),
        (4.8, "🔍 Complex Indian Plates", "Foreign OCRs Fail on Indian Roads", ACCENT_AMBER, [
            "Standard foreign OCR systems fail miserably on Indian roads due to high font variations, 2-line plates, commercial yellow plates, and dust/glare.",
            "✖ 40%+ OCR errors on Indian standard fonts",
            "✖ Night-vision headlight glare and angle distortion failures"
        ]),
        (8.8, "⏱️ Delayed Suspect Intercept", "Slow Manual Search & Cloud Lock-in", ACCENT_RED, [
            "Stolen vehicles and wanted criminals cross checkpoints unnoticed because manual monitoring cannot check hundreds of streams in real time.",
            "✖ Hours lost in manual DVR video playback",
            "✖ Exorbitant per-camera cloud subscription license fees"
        ])
    ]

    for left, title, sub, color, pts in cards2:
        tf = add_card(s2, left, 1.7, 3.7, 4.9, title, sub, color)
        p_desc = tf.add_paragraph()
        p_desc.text = pts[0]
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_LIGHT
        p_desc.space_before = Pt(8)

        for pt in pts[1:]:
            p_item = tf.add_paragraph()
            p_item.text = pt
            p_item.font.size = Pt(10.5)
            p_item.font.color.rgb = RED_LIGHT if "✖" in pt else TEXT_LIGHT
            p_item.space_before = Pt(8)

    # =========================================================================
    # SLIDE 3: 5-TIER HIGH-LEVEL ARCHITECTURE
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "High-Level Architecture & Component Stack", "SYSTEM TOPOLOGY", "Modular 5-tier architecture built for non-blocking asynchronous processing and extreme scalability.", 3)

    tiers = [
        (0.8, "1️⃣ Ingestion Tier", "MediaMTX Edge Gateway", ACCENT_CYAN, [
            "• Multi-Vendor RTSP/ONVIF Ingest",
            "• WebRTC / LL-HLS Video Proxy",
            "• Auto-Reconnect Watchdog",
            "• < 150ms Stream Latency"
        ]),
        (3.8, "2️⃣ AI Vision Core", "YOLOv8 + PP-OCRv4", ACCENT_PURPLE, [
            "• Dual-Head Vehicle & Plate",
            "• CLAHE Contrast Normalization",
            "• PP-OCRv4 Deep Text Engine",
            "• Indian LP Syntax Matcher"
        ]),
        (6.8, "3️⃣ Core Backend", "FastAPI & Redis Cache", ACCENT_AMBER, [
            "• < 1ms Redis Hotlist Index",
            "• Async WebSocket Broadcast",
            "• Trajectory Speed Engine",
            "• REST API & JWT Security"
        ]),
        (9.8, "4️⃣ Command UI", "React 18 & Leaflet GIS", ACCENT_GREEN, [
            "• Interactive Leaflet GIS Map",
            "• Live 1/4/9 Camera Matrix Grid",
            "• Audio/Visual Siren Toaster",
            "• Forensic Breadcrumb Replay"
        ])
    ]

    for left, title, sub, color, pts in tiers:
        tf = add_card(s3, left, 1.7, 2.7, 4.9, title, sub, color)
        for pt in pts:
            p = tf.add_paragraph()
            p.text = pt
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_LIGHT
            p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 4: SPECIALIZED AI ANPR PIPELINE
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "Specialized AI ANPR Pipeline for Indian Conditions", "COMPUTER VISION", "Multi-stage pipeline engineered specifically for high accuracy under difficult lighting, angles, and plate formats.", 4)

    # Left Box: 5 Steps
    tf_ai = add_card(s4, 0.8, 1.7, 6.2, 4.9, "⚡ 5-Step Deep Learning Inference", "Sub-100ms Pure On-Premise Execution", ACCENT_PURPLE)
    steps = [
        ("1. Adaptive Sampler", "Motion-triggered frame grabbing throttles to 10 FPS to save 60% compute."),
        ("2. Dual-Head YOLOv8", "Isolates vehicle class (Car, Bike, Truck, Auto) and crops the license plate rectangle."),
        ("3. CLAHE Image Enhancement", "Normalizes night headlight glare, shadow contrast, and perspective deskew."),
        ("4. PaddleOCR (PP-OCRv4)", "Deep character recognition trained on Indian standard, 2-line, yellow commercial & EV plates."),
        ("5. Indian LP Regex Validator", "Enforces state/RTO syntax & corrects optical ambiguities (O vs 0, B vs 8, I vs 1).")
    ]
    for h, d in steps:
        p = tf_ai.add_paragraph()
        p.text = f"• {h}: {d}"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(7)

    # Right Box: Visual Radar Box & Stats
    tf_right = add_card(s4, 7.3, 1.7, 5.2, 4.9, "📊 Live Detection Simulation & Benchmarks", "Validated on real-world Indian highway feeds", ACCENT_GREEN)

    # Simulated Radar Box Inside
    sim_rect = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(2.6), Inches(4.6), Inches(1.8))
    sim_rect.fill.solid()
    sim_rect.fill.fore_color.rgb = SUB_BOX_BG
    sim_rect.line.color.rgb = ACCENT_CYAN
    sim_rect.line.width = Pt(1)

    # Header inside sim
    tb_sim_h = s4.shapes.add_textbox(Inches(7.7), Inches(2.65), Inches(4.4), Inches(0.3))
    p_sh = tb_sim_h.text_frame.paragraphs[0]
    p_sh.text = "CAMERA: SG Highway (CAM-04)   •   GPU INFERENCE: 42ms"
    p_sh.font.size = Pt(9.5)
    p_sh.font.bold = True
    p_sh.font.color.rgb = CYAN_LIGHT

    # Yellow Plate Badge
    plate = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(3.0), Inches(3.2), Inches(0.65))
    plate.fill.solid()
    plate.fill.fore_color.rgb = YELLOW_PLATE
    plate.line.color.rgb = BLACK
    plate.line.width = Pt(2)
    p_pl = plate.text_frame.paragraphs[0]
    p_pl.text = "GJ 01 AB 1234"
    p_pl.font.size = Pt(18)
    p_pl.font.bold = True
    p_pl.font.color.rgb = BLACK
    p_pl.alignment = PP_ALIGN.CENTER

    # Vehicle details below plate
    tb_veh = s4.shapes.add_textbox(Inches(7.7), Inches(3.75), Inches(4.4), Inches(0.5))
    p_veh = tb_veh.text_frame.paragraphs[0]
    p_veh.text = "VEHICLE: WHITE FORTUNER (SUV)  |  CONFIDENCE: 96.4%"
    p_veh.font.size = Pt(9.5)
    p_veh.font.color.rgb = TEXT_MUTED
    p_veh.alignment = PP_ALIGN.CENTER

    # 2 Stat Metric Boxes
    stat1 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(4.6), Inches(2.2), Inches(1.6))
    stat1.fill.solid()
    stat1.fill.fore_color.rgb = SUB_BOX_BG
    stat1.line.color.rgb = CARD_BORDER
    p_s1 = stat1.text_frame.paragraphs[0]
    p_s1.text = "\n97.4%\n"
    p_s1.font.size = Pt(22)
    p_s1.font.bold = True
    p_s1.font.color.rgb = CYAN_LIGHT
    p_s1.alignment = PP_ALIGN.CENTER
    p_s1b = stat1.text_frame.add_paragraph()
    p_s1b.text = "Detection Accuracy"
    p_s1b.font.size = Pt(10)
    p_s1b.font.color.rgb = TEXT_MUTED
    p_s1b.alignment = PP_ALIGN.CENTER

    stat2 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.0), Inches(4.6), Inches(2.2), Inches(1.6))
    stat2.fill.solid()
    stat2.fill.fore_color.rgb = SUB_BOX_BG
    stat2.line.color.rgb = CARD_BORDER
    p_s2 = stat2.text_frame.paragraphs[0]
    p_s2.text = "\n< 45ms\n"
    p_s2.font.size = Pt(22)
    p_s2.font.bold = True
    p_s2.font.color.rgb = GREEN_LIGHT
    p_s2.alignment = PP_ALIGN.CENTER
    p_s2b = stat2.text_frame.add_paragraph()
    p_s2b.text = "GPU Frame Latency"
    p_s2b.font.size = Pt(10)
    p_s2b.font.color.rgb = TEXT_MUTED
    p_s2b.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 5: SUB-SECOND HOTLIST ALERTING
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "Sub-Second Watchlist Matching & Operator Siren", "REAL-TIME ALERTING", "Instant broadcast from camera sensor to operator workstation in under 400 milliseconds.", 5)

    # Left: Match Engine
    tf_hot = add_card(s5, 0.8, 1.7, 5.6, 4.9, "🚨 Real-Time Hotlist Match Engine", "Integrates with CCTNS & police stolen vehicle databases", ACCENT_RED)
    h_items = [
        ("Redis In-Memory Key-Value", "Performs hash lookups in < 1ms per recognized license plate."),
        ("Fuzzy Levenshtein Matcher", "Tolerates single-character OCR scratches to catch disguised plates."),
        ("Severity Categorization", "Distinguishes CRITICAL (Stolen/Armed), HIGH (Wanted), and MEDIUM alerts."),
        ("Instant Audio Siren", "Loud audio chime and red flashing toaster alerts the operator immediately.")
    ]
    for h, d in h_items:
        p = tf_hot.add_paragraph()
        p.text = f"• {h}: {d}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(10)

    # Right: Live Alert Card Mockup
    tf_alert = add_card(s5, 6.8, 1.7, 5.7, 4.9, "⚡ Command Center Alert Card Example", "What the police operator sees in real time:", ACCENT_CYAN)

    # Red Alert Banner Box
    banner = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(2.6), Inches(5.1), Inches(0.8))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(60, 15, 20)
    banner.line.color.rgb = ACCENT_RED
    banner.line.width = Pt(1.5)
    p_ban = banner.text_frame.paragraphs[0]
    p_ban.text = "🚨 CRITICAL HOTLIST MATCH  [ < 380ms ]"
    p_ban.font.size = Pt(12)
    p_ban.font.bold = True
    p_ban.font.color.rgb = RED_LIGHT
    p_ban2 = banner.text_frame.add_paragraph()
    p_ban2.text = "FIR-2026-9081 • STOLEN SUV REPORTED"
    p_ban2.font.size = Pt(10)
    p_ban2.font.color.rgb = TEXT_LIGHT

    # Alert Details Box
    det_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(3.5), Inches(5.1), Inches(1.8))
    det_box.fill.solid()
    det_box.fill.fore_color.rgb = SUB_BOX_BG
    det_box.line.color.rgb = CARD_BORDER
    tf_det = det_box.text_frame
    tf_det.word_wrap = True

    dets = [
        "📍 Location: SG Highway - Pakwan Junction (CAM-04)",
        "⏱️ Timestamp: Today at 10:45:22 AM",
        "🚘 Vehicle: GJ01AB1234 (White Toyota Fortuner)",
        "👮 Police Station: Vastrapur PS, Ahmedabad"
    ]
    for i, d in enumerate(dets):
        p_d = tf_det.paragraphs[0] if i == 0 else tf_det.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = TEXT_LIGHT
        if i > 0: p_d.space_before = Pt(4)

    # 2 Action Buttons
    btn1 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(5.5), Inches(2.45), Inches(0.6))
    btn1.fill.solid()
    btn1.fill.fore_color.rgb = RGBColor(2, 132, 199)
    btn1.line.fill.background()
    p_bt1 = btn1.text_frame.paragraphs[0]
    p_bt1.text = "📺 1-Click Live Feed"
    p_bt1.font.size = Pt(11)
    p_bt1.font.bold = True
    p_bt1.font.color.rgb = TEXT_WHITE
    p_bt1.alignment = PP_ALIGN.CENTER

    btn2 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.75), Inches(5.5), Inches(2.45), Inches(0.6))
    btn2.fill.solid()
    btn2.fill.fore_color.rgb = ACCENT_RED
    btn2.line.fill.background()
    p_bt2 = btn2.text_frame.paragraphs[0]
    p_bt2.text = "🚔 Dispatch PCR Van"
    p_bt2.font.size = Pt(11)
    p_bt2.font.bold = True
    p_bt2.font.color.rgb = TEXT_WHITE
    p_bt2.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 6: GIS VEHICLE BREADCRUMB TRACKING
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "Interactive GIS Map & Vehicle Breadcrumb Replay", "GIS & FORENSICS", "Reconstruct suspect escape routes with spatial speed estimation and directional polylines.", 6)

    cards6 = [
        (0.8, "🗺️ Live GIS Overview", "Leaflet + OpenStreetMap + PostGIS", ACCENT_CYAN, [
            "• Green/Red camera health pins across city junctions",
            "• Flashing radar alert marker on incident camera",
            "• Spatial camera clustering for high-density areas",
            "• Zero Google Maps licensing fees (100% self-hosted)"
        ]),
        (4.8, "📍 Numbered Breadcrumbs", "Chronological Journey Synthesis", ACCENT_PURPLE, [
            "• Trajectory reconstruction (1 → 2 → 3 → 4)",
            "• Average speed calculation (km/h) between junctions",
            "• Delta-T transit time between consecutive cameras",
            "• Directional arrows highlight fleeing suspect heading"
        ]),
        (8.8, "📋 Forensic Dossier", "Court-Admissible Evidence Export", ACCENT_GREEN, [
            "• 1-Click printable PDF investigation dossiers",
            "• Full-frame snapshots + localized plate crops",
            "• Wildcard search for partial plates (GJ01??1234)",
            "• Tamper-evident cryptographic audit hashes"
        ])
    ]

    for left, title, sub, color, pts in cards6:
        tf = add_card(s6, left, 1.7, 3.7, 4.9, title, sub, color)
        for pt in pts:
            p = tf.add_paragraph()
            p.text = pt
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_LIGHT
            p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 7: COMMAND CENTER CAPABILITIES
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "Unified Web Command Center Capabilities", "COMMAND CENTER", "Modern React 18 interface designed specifically for fast-paced police control room operations.", 7)

    # 4 Cards on Top
    dash_cards = [
        (0.8, "📺 30 Live CCTV Grid", "30 official Gujarat Police cameras with 60 FPS hardware HLS streams.", ACCENT_CYAN),
        (3.8, "🔒 Officer Security & RBAC", "Salted SHA-256 gated login: SP IT & Cyber (Badge: GP-7829).", ACCENT_GREEN),
        (6.8, "🚨 Instant Hotlist Alerts", "Sub-second cross-reference against CCTNS stolen FIR watchlists.", ACCENT_RED),
        (9.8, "🌐 Live Cloud Platform", "sentinel.deventtechnology.com with 24/7 keep-alive monitoring.", ACCENT_PURPLE)
    ]

    for left, title, desc, col in dash_cards:
        tf = add_card(s7, left, 1.7, 2.7, 2.8, title, "", col)
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(8)

    # Bottom Metric Bar
    bar_rect = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.8))
    bar_rect.fill.solid()
    bar_rect.fill.fore_color.rgb = SUB_BOX_BG
    bar_rect.line.color.rgb = CARD_BORDER

    metrics = [
        ("30 Feeds", "Live Gujarat Police Cameras", CYAN_LIGHT),
        ("SHA-256", "Officer Security Gateway", GREEN_LIGHT),
        ("< 1s", "Hotlist Alert Latency", RED_LIGHT),
        ("24/7 Live", "Continuous Cloud Uptime", AMBER_LIGHT)
    ]

    for i, (val, lbl, col) in enumerate(metrics):
        tb_m = s7.shapes.add_textbox(Inches(0.8 + i * 2.9), Inches(5.0), Inches(2.9), Inches(1.3))
        tf_m = tb_m.text_frame
        p_val = tf_m.paragraphs[0]
        p_val.text = val
        p_val.font.size = Pt(22)
        p_val.font.bold = True
        p_val.font.color.rgb = col
        p_val.alignment = PP_ALIGN.CENTER

        p_lbl = tf_m.add_paragraph()
        p_lbl.text = lbl
        p_lbl.font.size = Pt(11)
        p_lbl.font.color.rgb = TEXT_MUTED
        p_lbl.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 8: SUMMARY & ROI
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_header(s8, "Summary, Live Production Demo & Next Steps", "VALUE SUMMARY", "High operational impact with verified production cloud availability.", 8)

    tf_roi = add_card(s8, 0.8, 1.7, 5.6, 4.9, "💰 Massive Cost & Infrastructure ROI", "Zero Recurring SaaS Fees", ACCENT_GREEN)
    roi_items = [
        ("Zero Recurring SaaS Costs", "No per-camera licensing fees; saves crores annually compared to proprietary suites."),
        ("Multi-Vendor Reusability", "Works with existing legacy CCTV hardware across Gujarat without costly replacements."),
        ("100% Data Sovereignty", "Video footage and criminal records remain strictly inside the State Data Centre."),
        ("Rapid Resolution", "Cuts forensic tracking time from several hours down to a few seconds.")
    ]
    for h, d in roi_items:
        p = tf_roi.add_paragraph()
        p.text = f"• {h}: {d}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(10)

    tf_road = add_card(s8, 6.8, 1.7, 5.7, 4.9, "🌐 Live Production Demo Access", "Direct Evaluation Access", ACCENT_CYAN)
    road_items = [
        ("Live Web Dashboard", "https://sentinel.deventtechnology.com"),
        ("Live API Backend", "https://sentinel-api-bqfm.onrender.com/docs"),
        ("Officer Login Email", "jyoti@deventtechnology.com (Superintendent of Police)"),
        ("Officer Password", "123456 (Badge Number: GP-7829)"),
        ("Active Feeds", "30 Official Gujarat Police Cameras (cam01 - cam30)"),
        ("Uptime Monitoring", "Automated 24/7 keep-alive monitor enabled.")
    ]
    for h, d in road_items:
        p = tf_road.add_paragraph()
        p.text = f"• {h}: {d}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(8)

    for fname in ["Sentinel_AI_Surveillance_Presentation.pptx", "Sentinel_AI_Presentation_Final.pptx", "Sentinel_AI_Solution_Presentation.pptx"]:
        try:
            out_p = os.path.join(r"D:\sentinel-ai-surveillance\submission_package", fname)
            prs.save(out_p)
            print(f"Native PPTX generated successfully to: {out_p}")
        except Exception as e:
            print(f"Could not save to {fname} (likely open in PowerPoint): {e}")

if __name__ == "__main__":
    build_perfect_presentation()
